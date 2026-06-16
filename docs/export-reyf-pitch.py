#!/usr/bin/env python3
"""Exporta reyf-pitch.html a PowerPoint. La slide de Solución lleva el video incrustado."""

import asyncio
import sys
from pathlib import Path

DOCS = Path(__file__).resolve().parent
HTML = DOCS / "reyf-pitch.html"
VIDEO = DOCS / "reyf-demo.mp4"
OUT_PPTX = DOCS / "reyf-pitch.pptx"
TMP = DOCS / ".export-slides"
VIEWPORT = {"width": 1920, "height": 1080}
SOLUTION_SLIDE = 3  # 0-based — slide "Solución" con demo móvil


async def capture_slides() -> list[tuple[Path, dict | None]]:
    from playwright.async_api import async_playwright

    TMP.mkdir(exist_ok=True)
    for old in TMP.glob("slide-*.png"):
        old.unlink()

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport=VIEWPORT)
        await page.goto(HTML.as_uri(), wait_until="networkidle")
        await page.wait_for_timeout(1500)

        await page.add_style_tag(
            content="""
            .dots,.arrows,.hint{display:none!important}
            .rv,.lc .cell,.mini{opacity:1!important;transform:none!important}
            """
        )

        count = await page.evaluate("document.querySelectorAll('.slide').length")
        shots: list[tuple[Path, dict | None]] = []

        for i in range(count):
            video_box = None
            if i == SOLUTION_SLIDE:
                video_box = await page.evaluate(
                    """() => {
                      const r = document.querySelector('.phone-screen')?.getBoundingClientRect();
                      if (!r) return null;
                      const v = document.querySelector('.phone-screen video');
                      if (v) v.style.visibility = 'hidden';
                      return {x: r.x, y: r.y, width: r.width, height: r.height};
                    }"""
                )

            await page.evaluate(
                """(idx) => {
                  const slides = [...document.querySelectorAll('.slide')];
                  const track = document.getElementById('track');
                  track.style.transform = `translateX(${-idx * 100}vw)`;
                  slides.forEach((s, n) => s.classList.toggle('active', n === idx));
                  slides[idx].querySelectorAll('[data-count]').forEach((el) => {
                    if (el.dataset.done) return;
                    el.dataset.done = '1';
                    const target = parseFloat(el.dataset.count);
                    const dec = parseInt(el.dataset.dec || '0', 10);
                    const pre = el.dataset.pre || '';
                    const suf = el.dataset.suf || '';
                    el.textContent = pre + target.toFixed(dec) + suf;
                  });
                }""",
                i,
            )
            await page.wait_for_timeout(1600 if i in (1, SOLUTION_SLIDE) else 1000)

            path = TMP / f"slide-{i:02d}.png"
            await page.screenshot(path=str(path))
            shots.append((path, video_box))
            print(f"  capturada slide {i + 1}/{count}" + (" + video" if video_box else ""))

            if i == SOLUTION_SLIDE:
                await page.evaluate(
                    "() => { const v = document.querySelector('.phone-screen video'); if (v) v.style.visibility = ''; }"
                )

        await browser.close()
    return shots


def px_to_inches(px: float, axis: str) -> float:
    from pptx.util import Inches

    w, h = VIEWPORT["width"], VIEWPORT["height"]
    slide_w, slide_h = 13.333, 7.5
    if axis == "x":
        return Inches(px * slide_w / w)
    return Inches(px * slide_h / h)


def build_pptx(slides: list[tuple[Path, dict | None]]) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    if not VIDEO.exists():
        raise SystemExit(f"Falta el video {VIDEO}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for img, video_box in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
        if video_box:
            slide.shapes.add_movie(
                str(VIDEO),
                px_to_inches(video_box["x"], "x"),
                px_to_inches(video_box["y"], "y"),
                px_to_inches(video_box["width"], "x"),
                px_to_inches(video_box["height"], "y"),
                mime_type="video/mp4",
            )
            print("  video incrustado en slide Solución")

    prs.save(OUT_PPTX)


async def main() -> None:
    if not HTML.exists():
        raise SystemExit(f"No se encontró {HTML}")

    print("Capturando slides…")
    slides = await capture_slides()
    print("Generando PowerPoint…")
    build_pptx(slides)
    print(f"Listo: {OUT_PPTX}")
    print("Tip: en PowerPoint, entra a la slide Solución y pulsa Play en el video.")


if __name__ == "__main__":
    try:
        asyncio.run(main())
    except ImportError as exc:
        print("Faltan dependencias. Instala con:", file=sys.stderr)
        print("  cd Seyf2/docs && python3 -m venv .venv", file=sys.stderr)
        print("  .venv/bin/pip install python-pptx playwright", file=sys.stderr)
        print("  .venv/bin/playwright install chromium", file=sys.stderr)
        raise SystemExit(1) from exc
